"""
文档加载器 —— 多格式文档文本提取

学习要点：
- 了解不同文档格式（PDF、Word、Excel、PPT）的解析方式
- 理解 RAG 第一步：将非结构化文档转换为纯文本
"""

import os
import logging
from io import StringIO


def extract_text(filepath):
    """
    从文件中提取纯文本内容

    支持格式：PDF / Word / Excel / PPT / 纯文本 / Markdown

    Args:
        filepath: 文件路径

    Returns:
        提取的文本内容字符串
    """
    file_ext = os.path.splitext(filepath)[1].lower()

    if file_ext == '.pdf':
        from pdfminer.high_level import extract_text_to_fp
        output = StringIO()
        with open(filepath, 'rb') as file:
            extract_text_to_fp(file, output)
        return output.getvalue()

    elif file_ext in ['.txt', '.md']:
        with open(filepath, 'r', encoding='utf-8') as file:
            return file.read()

    elif file_ext == '.docx':
        try:
            from docx import Document
            doc = Document(filepath)
            return "\n".join([para.text for para in doc.paragraphs])
        except ImportError:
            logging.error("处理Word文档需要安装python-docx库")
            return ""

    elif file_ext in ['.xlsx', '.xls']:
        try:
            import pandas as pd
            text = ""
            xl = pd.ExcelFile(filepath)
            for sheet_name in xl.sheet_names:
                df = xl.parse(sheet_name)
                text += f"工作表: {sheet_name}\n"
                text += df.to_string(index=False) + "\n\n"
            return text
        except ImportError:
            logging.error("处理Excel文件需要安装pandas库")
            return ""

    elif file_ext == '.pptx':
        try:
            from pptx import Presentation
            prs = Presentation(filepath)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text
        except ImportError:
            logging.error("处理PPT文件需要安装python-pptx库")
            return ""

    else:
        logging.warning(f"不支持的文件格式: {file_ext}")
        return ""
